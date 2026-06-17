#!/usr/bin/env python3
"""DeckForge local harness utilities.

This script provides deterministic glue for the DeckForge skills. It does not
replace Playwright MCP or frontend-design; it gives Codex repeatable local
commands for the file and render stages around those tools.
"""

from __future__ import annotations

import argparse
import json
import platform
import shutil
import subprocess
import sys
from pathlib import Path

from PIL import Image, ImageDraw, ImageOps, ImageStat
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt


SLIDE_W = 13.333333
SLIDE_H = 7.5


def fail(message: str) -> None:
    raise SystemExit(f"error: {message}")


def which(name: str) -> str:
    exe = shutil.which(name)
    if not exe:
        fail(f"missing required executable: {name}")
    return exe


def cmd_init(args: argparse.Namespace) -> None:
    root = Path(args.project).resolve()
    for child in ["captures", "slides", "renders", "pptx", "qa", "harness"]:
        (root / child).mkdir(parents=True, exist_ok=True)
    sources = root / "captures" / "sources.json"
    if not sources.exists():
        sources.write_text(json.dumps({"sources": []}, indent=2) + "\n", encoding="utf-8")
    print(root)


def cmd_doctor(args: argparse.Namespace) -> None:
    system = platform.system()
    checks = {
        "platform": system,
        "python": sys.executable,
        "node": shutil.which("node"),
        "soffice": shutil.which("soffice") or shutil.which("libreoffice"),
        "pdftoppm": shutil.which("pdftoppm"),
        "pywin32": False,
        "recommended_backend": "python-pptx + LibreOffice",
    }
    if system == "Windows":
        try:
            import win32com.client  # type: ignore

            checks["pywin32"] = True
            checks["recommended_backend"] = "WPS/MS Office COM for live Office automation; python-pptx for file edits"
        except Exception:
            checks["recommended_backend"] = "python-pptx; install pywin32 + WPS/MS Office for COM automation"
    elif system == "Darwin":
        checks["recommended_backend"] = "python-pptx + LibreOffice; WPS COM is unavailable on macOS"
    print(json.dumps(checks, ensure_ascii=False, indent=2))


def cmd_html_to_png(args: argparse.Namespace) -> None:
    html = Path(args.html).resolve()
    if not html.exists():
        fail(f"missing HTML file: {html}")
    out_dir = Path(args.out_dir).resolve()
    out_dir.mkdir(parents=True, exist_ok=True)

    node = which("node")
    js = f"""
const {{ chromium }} = await import('playwright');
const browser = await chromium.launch({{ headless: true }});
const page = await browser.newPage({{ viewport: {{ width: {args.width}, height: {args.height} }}, deviceScaleFactor: 1 }});
await page.goto('file://{html}', {{ waitUntil: 'networkidle' }});
const slides = await page.locator('[data-slide]').count();
if (slides === 0) {{
  await page.screenshot({{ path: '{out_dir / "slide-01.png"}', fullPage: false }});
}} else {{
  for (let i = 0; i < slides; i++) {{
    const el = page.locator('[data-slide]').nth(i);
    await el.screenshot({{ path: `{out_dir}/slide-${{String(i + 1).padStart(2, '0')}}.png` }});
  }}
}}
await browser.close();
"""
    subprocess.run([node, "--input-type=module", "-e", js], check=True)
    print(f"rendered_dir: {out_dir}")


def image_files(path: Path, pattern: str) -> list[Path]:
    return [
        p for p in sorted(path.glob(pattern))
        if p.suffix.lower() in {".png", ".jpg", ".jpeg"}
    ]


def cmd_images_to_pptx(args: argparse.Namespace) -> None:
    input_dir = Path(args.input_dir).resolve()
    images = image_files(input_dir, args.pattern)
    if not images:
        fail(f"no images found in {input_dir} matching {args.pattern}")
    output = Path(args.output).resolve()
    output.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]
    sw = float(prs.slide_width)
    sh = float(prs.slide_height)

    for idx, image in enumerate(images, start=1):
        slide = prs.slides.add_slide(blank)
        with Image.open(image) as im:
            iw, ih = im.size
        ratio = iw / ih
        slide_ratio = sw / sh
        if ratio > slide_ratio:
            h = sh
            w = h * ratio
            x = (sw - w) / 2
            y = 0
        else:
            w = sw
            h = w / ratio
            x = 0
            y = (sh - h) / 2
        pic = slide.shapes.add_picture(str(image), x, y, width=int(w), height=int(h))
        pic.name = f"DeckForge full-slide image {idx:02d}"

    prs.save(output)
    print(f"pptx: {output}")
    print(f"slides: {len(images)}")


def hex_color(value: str, default: str = "000000") -> str:
    text = str(value or default).strip().replace("#", "")
    if len(text) == 3:
        text = "".join(ch * 2 for ch in text)
    try:
        int(text, 16)
    except ValueError:
        return default
    return text.upper() if len(text) == 6 else default


def px_scale(canvas: dict) -> tuple[float, float]:
    w = float(canvas.get("w", 1280))
    h = float(canvas.get("h", 720))
    if w <= 0 or h <= 0:
        fail("canvas.w and canvas.h must be positive")
    return SLIDE_W / w, SLIDE_H / h


def add_schema_text(slide, el: dict, sx: float, sy: float) -> None:
    shape = slide.shapes.add_textbox(
        Inches(float(el.get("x", 0)) * sx),
        Inches(float(el.get("y", 0)) * sy),
        Inches(float(el.get("w", 200)) * sx),
        Inches(float(el.get("h", 50)) * sy),
    )
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = bool(el.get("wrap", True))
    if bool(el.get("fit", True)):
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = str(el.get("text", ""))
    p.font.size = Pt(float(el.get("fs", 24)))
    p.font.bold = bool(el.get("bold", False))
    p.font.name = str(el.get("font", "Arial"))
    p.font.color.rgb = RGBColor.from_string(hex_color(el.get("color", "111827")))
    align = str(el.get("align", "left")).lower()
    if align in {"center", "2"}:
        p.alignment = 2
    elif align in {"right", "3"}:
        p.alignment = 3


def add_schema_rect(slide, el: dict, sx: float, sy: float, rounded: bool = False) -> None:
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        kind,
        Inches(float(el.get("x", 0)) * sx),
        Inches(float(el.get("y", 0)) * sy),
        Inches(float(el.get("w", 100)) * sx),
        Inches(float(el.get("h", 40)) * sy),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(hex_color(el.get("fill", el.get("color", "FFFFFF")), "FFFFFF"))
    line_color = el.get("line")
    if line_color:
        shape.line.color.rgb = RGBColor.from_string(hex_color(line_color))
        shape.line.width = Pt(float(el.get("line_width", 1)))
    else:
        shape.line.fill.background()


def add_schema_image(slide, el: dict, sx: float, sy: float, base_dir: Path) -> None:
    image = Path(str(el.get("file", "")))
    if not image.is_absolute():
        image = base_dir / image
    if not image.exists():
        fail(f"schema image not found: {image}")
    slide.shapes.add_picture(
        str(image),
        Inches(float(el.get("x", 0)) * sx),
        Inches(float(el.get("y", 0)) * sy),
        width=Inches(float(el.get("w", 100)) * sx),
        height=Inches(float(el.get("h", 100)) * sy),
    )


def add_schema_line(slide, el: dict, sx: float, sy: float) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(float(el.get("x", 0)) * sx),
        Inches(float(el.get("y", 0)) * sy),
        Inches(float(el.get("w", 100)) * sx),
        Inches(max(float(el.get("h", 2)) * sy, 0.01)),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(hex_color(el.get("color", "111827")))
    shape.line.fill.background()


def cmd_schema_to_pptx(args: argparse.Namespace) -> None:
    schema_path = Path(args.schema).resolve()
    if not schema_path.exists():
        fail(f"missing schema file: {schema_path}")
    data = json.loads(schema_path.read_text(encoding="utf-8"))
    output = Path(args.output).resolve()
    output.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]
    canvas = data.get("canvas", {"w": 1280, "h": 720})
    sx, sy = px_scale(canvas)
    base_dir = schema_path.parent

    for slide_data in data.get("slides", []):
        slide = prs.slides.add_slide(blank)
        bg = slide_data.get("background")
        if bg:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor.from_string(hex_color(bg, "FFFFFF"))
        for el in slide_data.get("elements", []):
            etype = str(el.get("type", "text"))
            if etype == "text":
                add_schema_text(slide, el, sx, sy)
            elif etype in {"rect", "shape"}:
                add_schema_rect(slide, el, sx, sy, rounded=False)
            elif etype in {"rrect", "card"}:
                add_schema_rect(slide, el, sx, sy, rounded=True)
            elif etype == "image":
                add_schema_image(slide, el, sx, sy, base_dir)
            elif etype == "line":
                add_schema_line(slide, el, sx, sy)
            else:
                fail(f"unsupported schema element type: {etype}")

    prs.save(output)
    print(f"pptx: {output}")
    print(f"slides: {len(data.get('slides', []))}")


def cmd_render_pptx(args: argparse.Namespace) -> None:
    pptx = Path(args.pptx).resolve()
    if not pptx.exists():
        fail(f"missing PPTX file: {pptx}")
    out_dir = Path(args.out_dir).resolve()
    out_dir.mkdir(parents=True, exist_ok=True)

    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        fail("missing LibreOffice executable: soffice/libreoffice")
    pdftoppm = which("pdftoppm")

    subprocess.run([soffice, "--headless", "--convert-to", "pdf", "--outdir", str(out_dir), str(pptx)], check=True)
    pdf = out_dir / f"{pptx.stem}.pdf"
    if not pdf.exists():
        fail(f"expected PDF was not created: {pdf}")
    subprocess.run([pdftoppm, "-png", "-r", str(args.dpi), str(pdf), str(out_dir / "slide")], check=True)
    rendered = sorted(out_dir.glob("slide-*.png"))
    print(f"pdf: {pdf}")
    print(f"rendered: {len(rendered)}")


def cmd_qa(args: argparse.Namespace) -> None:
    render_dir = Path(args.render_dir).resolve()
    images = image_files(render_dir, args.pattern)
    if not images:
        fail(f"no render images found in {render_dir}")
    failures: list[str] = []
    for image in images:
        im = Image.open(image).convert("RGB")
        stat = ImageStat.Stat(im)
        variance = sum(stat.var) / 3
        bbox = ImageOps.invert(im.convert("L")).getbbox()
        print(f"{image.name}: size={im.size} variance={variance:.2f} nonwhite_bbox={bbox}")
        if bbox is None or variance < args.min_variance:
            failures.append(image.name)
    if args.contact_sheet:
        make_contact_sheet(images, Path(args.contact_sheet).resolve())
    if failures:
        print("failed:", ", ".join(failures))
        sys.exit(1)
    print(f"qa_ok: {len(images)}")


def make_contact_sheet(images: list[Path], output: Path) -> None:
    thumbs = []
    for image in images:
        im = Image.open(image).convert("RGB")
        im.thumbnail((320, 180))
        canvas = Image.new("RGB", (320, 205), "white")
        canvas.paste(im, ((320 - im.width) // 2, 0))
        ImageDraw.Draw(canvas).text((8, 184), image.name, fill=(0, 0, 0))
        thumbs.append(canvas)
    cols = 2
    rows = (len(thumbs) + cols - 1) // cols
    sheet = Image.new("RGB", (cols * 320, rows * 205), "#eeeeee")
    for idx, thumb in enumerate(thumbs):
        sheet.paste(thumb, ((idx % cols) * 320, (idx // cols) * 205))
    output.parent.mkdir(parents=True, exist_ok=True)
    sheet.save(output)
    print(f"contact_sheet: {output}")


def parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(description="DeckForge local harness")
    sub = p.add_subparsers(dest="cmd", required=True)

    init = sub.add_parser("init")
    init.add_argument("--project", default="deckforge")
    init.set_defaults(func=cmd_init)

    doctor = sub.add_parser("doctor")
    doctor.set_defaults(func=cmd_doctor)

    html = sub.add_parser("html-to-png")
    html.add_argument("--html", required=True)
    html.add_argument("--out-dir", required=True)
    html.add_argument("--width", type=int, default=1280)
    html.add_argument("--height", type=int, default=720)
    html.set_defaults(func=cmd_html_to_png)

    imgs = sub.add_parser("images-to-pptx")
    imgs.add_argument("--input-dir", required=True)
    imgs.add_argument("--pattern", default="slide-*.png")
    imgs.add_argument("--output", required=True)
    imgs.set_defaults(func=cmd_images_to_pptx)

    schema = sub.add_parser("schema-to-pptx")
    schema.add_argument("--schema", required=True)
    schema.add_argument("--output", required=True)
    schema.set_defaults(func=cmd_schema_to_pptx)

    render = sub.add_parser("render-pptx")
    render.add_argument("--pptx", required=True)
    render.add_argument("--out-dir", required=True)
    render.add_argument("--dpi", type=int, default=120)
    render.set_defaults(func=cmd_render_pptx)

    qa = sub.add_parser("qa")
    qa.add_argument("--render-dir", required=True)
    qa.add_argument("--pattern", default="slide-*.png")
    qa.add_argument("--min-variance", type=float, default=1.0)
    qa.add_argument("--contact-sheet")
    qa.set_defaults(func=cmd_qa)

    return p


def main() -> None:
    args = parser().parse_args()
    args.func(args)


if __name__ == "__main__":
    main()
